import os
import numpy as np
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from prettytable import PrettyTable as PT
from PIL import Image
from matplotlib import pyplot as plt
import time
import random
########### Own Module #########
from ghost_line import double_line
from worker import save_2d_array_as_png


def get_random_number():
    # Get the current time in microseconds
    microseconds = int(time.time() * 1e6)

    # Seed the random number generator with the microseconds
    random.seed(microseconds)

    # Generate a random number
    return random.random()


PNG_FILE1 = 'png/' + str(get_random_number()) + '.png'
PNG_FILE2 = 'png/' + str(get_random_number()) + '.png'


def ppt_pngs(filename, css):
    image = Image.open(filename)
    # Convert the image to grayscale
    grayscale_image = image.convert('L')
    # Convert the grayscale image to a 2D NumPy array
    a = np.array(grayscale_image)
    save_2d_array_as_png(a[114:1050, 1032:1811], PNG_FILE2)

    fig, axs = plt.subplots(3, 3)
    if len(css) == 36:
        for i in range(9):
            for j in range(4):
                l = css[i * 4 + j]
                ax = axs[int(i / 3)][i % 3]
                ax.plot(l, alpha=.5)
    plt.savefig(PNG_FILE1)
    plt.close()
    return


def get_png_files(directory):
    # Get a list of all PNG files in the specified directory
    png_files = glob.glob(os.path.join(directory, '*.png'))
    return png_files


def add_slide(prs, txt):
    # Add a blank slide layout
    slide_layout = prs.slide_layouts[6]  # 6 is for a blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Add text box on the left
    left = Inches(0.5)
    top = Inches(.5)
    width = Inches(7)
    height = Inches(4)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = txt
    p.font.size = Pt(12)  # Reduced from 28
    p.font.name = 'Courier New'

    # Add the plot image on the right
    img_path = PNG_FILE1
    left = Inches(.5)
    top = Inches(4)
    height = Inches(4.5)
    slide.shapes.add_picture(img_path, left, top, height=height)
    # Add the plot image on the right
    img_path = PNG_FILE2
    left = Inches(8)
    top = Inches(.5)
    height = Inches(8)
    slide.shapes.add_picture(img_path, left, top, height=height)


def new_deck():
    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs


def make_table(lvls, data):
    t = PT()
    t.field_names = ["Pos.", "Off Ctr", "AD Norm.",
                     "D1", "D2", "Score", "Lvl"]
    loc = ["Top Left", "Top Ctr", "Top Right",
           "Ctr Left", "Ctr Ctr", "Ctr Right",
           "Btm Left", "Btm Ctr", "Btm Right"]
    for i in range(9):
        t.add_row([loc[i], "%.2f" % data[5][i],
                   "%.2f" % data[1][i], "%.2f" % data[3][i],
                   "%.2f" % data[4][i], "%.2f" % data[0][i],
                   "%d" % lvls[i]])
    return str(t)


def main2(listfile):
    def read_file_to_list(filename):
        with open(filename, 'r') as file:
            # Read all lines and strip whitespace
            lines = [line.strip() for line in file]
        return lines

    files = read_file_to_list(listfile)
    record = open('png/' + listfile + '.txt', 'w')
    print(files)
    prs = new_deck()
    for file in files:
        try:
            res = double_line(file)
            ppt_pngs(file, res["cross sections"])
            add_slide(prs, file + "\n" +
                      make_table(res["lvl"], res["data"]))
            output = file + "," + str(
                res['lvl']) + str(res['data'])
            print(output)
            record.write(output + "\n")
        except BaseException:
            pass
    prs.save('png/' + listfile + '.pptx')
    record.close()


if __name__ == "__main__":
    listfile = input("")
    main2(listfile)
    #  main3()
