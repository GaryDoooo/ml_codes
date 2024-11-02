import os
import subprocess
import cv2
import matplotlib.pyplot as plt
import matplotlib.colors as colors
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import uuid
##### Own Modules #####
from double_line_single import double_line

# Function to determine text color based on background luminance


def get_text_color(value, vmin, vmax, cmap):
    norm = plt.Normalize(vmin, vmax)
    rgba = cmap(norm(value))
    luminance = 0.299 * rgba[0] + 0.587 * rgba[1] + 0.114 * rgba[2]
    return 'black' if luminance > 0.5 else 'white'

# Create a custom colormap


def create_custom_ylord():
    # Get the YlOrRd colormap
    ylord = plt.cm.YlOrRd(np.linspace(0, 1, 256))

    # Create a new colormap with white at the beginning
    custom_ylord = np.vstack((np.array([1, 1, 1, 1]), ylord))
    return colors.ListedColormap(custom_ylord)


def get_tiff_filenames(directory):
    tiff_files = []
    for filename in os.listdir(directory):
        if filename.lower().endswith(('.tif', '.tiff')):
            tiff_files.append(filename)
    return tiff_files


def plot_121_heatmap(res_map, png_filename):
    fig, ax = plt.subplots(figsize=(7, 7))
    ax.set_aspect('equal')
    y = x = [i for i in range(1, 12)]
    # Create a modified YlOrRd colormap with white as the minimum
    cmap = plt.cm.YlOrRd.copy()
    cmap.set_under('white')

    # Create the custom colormap
    cmap = create_custom_ylord()
    ax.pcolormesh(x, y, res_map, shading='auto', cmap=cmap, edgecolors='grey',
                  linewidths=0.5, vmin=1, vmax=4)

    # Add value annotations
    max_value = 4  # max(max(row) for row in res_map)
    min_value = 1  # min(min(row) for row in res_map)
    for i in range(len(y)):
        for j in range(len(x)):
            ax.text(x[j], y[i], f'{res_map[i][j]}',
                    ha='center', va='center', color=get_text_color(res_map[i][j],
                                                                   min_value, max_value, cmap))

    # Get the current y-axis limits
    y_min, y_max = ax.get_ylim()
    ax.set_ylim(y_max, y_min)
    ax.set_xticks(x)
    ax.set_xticklabels([f'C{i}' for i in x])
    ax.set_yticks(y)
    ax.set_yticklabels([f'R{i}' for i in y])
    # ax.set_title('2D Heatmap')
    # Move x-axis to the top
    ax.xaxis.tick_top()
    ax.xaxis.set_label_position('top')

    plt.savefig(png_filename)
    plt.close()


def cropNenhance(image, left, right, top, bottom):
    cropped_image = image[top:bottom + 1, left:right + 1]
    # Normalize the pixel values
    min_val = np.min(cropped_image)
    max_val = np.max(cropped_image)
    # Scale to [0, 255]
    enhanced_image = (255 * (cropped_image - min_val) / (max_val - min_val)).astype(np.uint8)
    return enhanced_image


def get_lvl_map(output_str, filename):
    file1 = str(uuid.uuid4()) + ".png"
    file2 = str(uuid.uuid4()) + ".png"
    xhs = data = eval(output_str.replace("lines parallel", ""))
    print(len(data), filename)
    xhs.sort(key=lambda x: x[1])
    xh2d = [xhs[i * 11:(i + 1) * 11] for i in range(11)]
    for i in xh2d:
        i.sort(key=lambda x: x[0])

    # Read the TIFF file with IMREAD_ANYDEPTH flag to handle 32-bit depth
    image = cv2.imread(filename, cv2.IMREAD_ANYDEPTH)

    # Convert to float if necessary (OpenCV reads it as uint16 or similar)
    image = image.astype('float32')

    res_map = [[0] * 11 for i in range(11)]
    for x in range(11):
        for y in range(11):
            [x0, y0] = xh2d[y][x]
            res = double_line(cropNenhance(image, x0 - 80, x0 + 80, y0 - 80, y0 + 80), x=80, y=80)
            res_map[y][x] = res['lvl']

    plot_121_heatmap(res_map, file1)

#  filename='2Y0Y3B1XSDG8Q002H_G.tiff'
    height, width = image.shape[:2]
    gridleft = xh2d[0][0][0]
    gridtop = xh2d[0][0][1]
    gridright = xh2d[10][10][0]
    gridbottom = xh2d[10][10][1]
    #  print(gridtop, gridbottom)

    cropleft = int(max(gridleft - (gridright - gridleft) / 10, 0))
    cropright = int(min(gridright + (gridright - gridleft) / 10, int(width)))

    croptop = int(max(gridtop - (gridbottom - gridtop) / 10, 0))
    cropbottom = int(min(gridbottom + (gridbottom - gridtop) / 10, int(height)))
# Create the plot with custom size
    fig, ax = plt.subplots(figsize=(7, 7))

# Display the image
#  im = ax.imshow(data, cmap='viridis', origin='lower')
    ax.imshow(cropNenhance(image, cropleft, cropright, croptop, cropbottom), cmap='gray')

# Overlay scatter points
    x = [i[0] - cropleft for i in xhs]
    y = [i[1] - croptop for i in xhs]
    ax.scatter(x, y, c='red', s=2)

    #  plt.figure(figsize=(7, 7))
    #  plt.imshow(cropNenhance(image, cropleft, cropright, croptop, cropbottom), cmap='gray')
    # plot the cross center detected
    #  for [x, y] in xhs:
    #      plt.plot(x - cropleft, y - cropright, 'ro')

    plt.title(filename.split('/')[-1])
    # Adjust the layout
    plt.tight_layout()
    plt.savefig(file2)
    plt.close()
    return file1, file2


def main(dir_name):
    directory = dir_name
    tiff_filenames = [os.path.join(dir_name, i) for i in get_tiff_filenames(directory)]
    print(tiff_filenames)

    # Add to PowerPoint
    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Define the command to be executed
    for filename in tiff_filenames:
        command = 'echo "' + filename + '" | ./find121x'
        # Run the command using subprocess.run
        result = subprocess.run(command, shell=True, capture_output=True, text=True)
        # Capture the output
        output = result.stdout
        # Print or use the captured output
        #  print(output)
        file1, file2 = get_lvl_map(output, filename)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(file2, Inches(0.5), Inches(1))
        slide.shapes.add_picture(file1, Inches(8), Inches(1))
        os.remove(file1)
        os.remove(file2)

    prs.save('pngdd.pptx')
    return


if __name__ == "__main__":
    dir_name = input("")
    main(dir_name)
