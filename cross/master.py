import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
########### Own Module #########
from worker import double_line


def get_png_files(directory):
    # Get a list of all PNG files in the specified directory
    png_files = glob.glob(os.path.join(directory, '*.png'))
    return png_files


def add_slide(prs, txt):
    # Add a blank slide layout
    slide_layout = prs.slide_layouts[6]  # 6 is for a blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Add text box on the left
    left = Inches(0.5)
    top = Inches(.5)
    width = Inches(7)
    height = Inches(4)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = txt
    p.font.size = Pt(12)  # Reduced from 28
    p.font.name = 'Courier New'

    # Add the plot image on the right
    img_path = 'png/1.png'
    left = Inches(.5)
    top = Inches(4)
    height = Inches(4.5)
    slide.shapes.add_picture(img_path, left, top, height=height)
    # Add the plot image on the right
    img_path = 'png/2.png'
    left = Inches(8)
    top = Inches(.5)
    height = Inches(8)
    slide.shapes.add_picture(img_path, left, top, height=height)


def new_deck():
    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs


def main():

    directory_path = './png'
    png_files = get_png_files(directory_path)
    # Create a presentation object

    for cnt, file in enumerate(png_files):
        if cnt == 0:
            prs = new_deck()
        elif cnt % 10 == 0:
            prs.save('png/ppt' + str(cnt) + '.pptx')
            prs = new_deck()
        try:
            res = double_line(file)
            add_slide(prs, res['str'])
            print(file + "," + str(
                res['lvl']) + str(res['data']))
        except BaseException:
            pass

    # Save the presentation
    prs.save('png/ppt' + '.pptx')
    #  prs.save('png/ppt.pptx')


def main2(listfile):
    def read_file_to_list(filename):
        with open(filename, 'r') as file:
            # Read all lines and strip whitespace
            lines = [line.strip() for line in file]
        return lines

    files = read_file_to_list(listfile)
    record = open('png/' + listfile + '.txt', 'w')
    print(files)
    prs = new_deck()
    for file in files:
        try:
            res = double_line(file)
            add_slide(prs, res['str'])
            output = file + "," + str(
                res['lvl']) + str(res['data'])
            print(output)
            record.write(output + "\n")
        except BaseException:
            pass
    prs.save('png/' + listfile + '.pptx')
    record.close()


if __name__ == "__main__":
    listfile = input("")
    main2(listfile)
